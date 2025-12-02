import os
import re
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from sqlalchemy import text
from extractor.cleaner import split_text_smartly
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

SLIDE_CHAR_LIMIT = 850
OUTPUT_FOLDER = "FDA_PPTs"

os.makedirs(OUTPUT_FOLDER, exist_ok=True)


def clean_filename(name):
    return re.sub(r'[\\/*?:"<>|]', "", name)[:50]


def clear_all_text(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()


def get_title_shape(slide):
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title
    # fallback
    return slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.2))


def get_body_shape(slide):
    # try placeholder index 1
    try:
        ph = slide.placeholders[1]
        if ph.has_text_frame:
            return ph
    except:
        pass

    # try any other placeholder
    for ph in slide.placeholders:
        if getattr(ph, "has_text_frame", False) and ph != slide.shapes.title:
            return ph

    # create a textbox fallback
    return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))


def copy_slide(prs, slide):
    """Create a new slide with the same layout."""
    layout = slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    return new_slide

def add_thank_you_slide(prs, message="Thank You", layout_index=3):
    """
    Adds a Thank You slide using an existing template layout.
    Ensures proper design and formatting.
    """
    # Use the template's own layout
    try:
        layout = prs.slide_layouts[layout_index]
    except:
        layout = prs.slide_layouts[-1]   # fallback to last layout

    slide = prs.slides.add_slide(layout)

    # Find title placeholder first
    title_shape = slide.shapes.title
    if title_shape and title_shape.has_text_frame:
        tf = title_shape.text_frame
        tf.text = message
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 128)
        p.alignment = 1
        return slide

    # If template has no title placeholder, create a text box
    textbox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    tf = textbox.text_frame
    tf.text = message
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 128)
    p.alignment = 1

    return slide


def add_watermark(slide, text="."):
    """Adds a rotated watermark to the given slide."""
    left = Inches(2.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(42)
    font.bold = True
    font.fill.transparency = 0.5
    font.color.rgb = RGBColor(200, 200, 200)   # Light gray watermark color
    textbox.rotation = 315                     # Diagonal watermark (45° upward left)

    # Send watermark behind all content (z-order)
    textbox.shadow.inherit = False
    return textbox


def parse_summary_text(summary_text: str) -> OrderedDict:
    """
    Converts the LLM summary string (with headings + '- bullet' lines)
    into an ordered dictionary of {heading: [bullets]}.
    """
    sections = OrderedDict()
    current_heading = None

    for line in summary_text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("-"):
            bullet = stripped.lstrip("- ").strip()
            if not bullet:
                continue
            if current_heading is None:
                current_heading = "Summary"
                sections.setdefault(current_heading, [])
            sections.setdefault(current_heading, []).append(bullet)
        else:
            current_heading = stripped
            sections.setdefault(current_heading, [])

    return sections

def wrap_url(url, max_len=60):
    return "\n".join(url[i:i+max_len] for i in range(0, len(url), max_len))




def create_individual_ppt(item):
    prs = Presentation("ppt/template.pptx")

    summary_source = item.get("Summary", "")
    if isinstance(summary_source, dict):
        summary_dict = summary_source
    elif isinstance(summary_source, str):
        summary_dict = parse_summary_text(summary_source)
    else:
        raise TypeError("item['Summary'] must be either a dict or a string.")

    # ---------------------------
    # SLIDE 1 (EDIT TITLE SLIDE)
    # ---------------------------
    
    slide1 = prs.slides[0]
    clear_all_text(slide1)
    add_watermark(slide1)


    text_shapes = [s for s in slide1.shapes if s.has_text_frame]

    if len(text_shapes) >= 1:
        title = text_shapes[0].text_frame
        title.text = item["Company"]
        p = title.paragraphs[0]
        p.font.size = Pt(30)
        p.font.name = "times new roman"
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255) # white

    if len(text_shapes) >= 2:
        subtitle = text_shapes[1].text_frame
        subtitle.text = f"FDA Warning Letter Report\nDate Posted: {item['Date Posted']}"

    # ---------------------------
    # SUMMARY TEMPLATE (Slide 2)
    # ---------------------------
    
    summary_template = prs.slides[1]

    # ---------------------------
    # BUILD EACH SECTION SLIDE
    # ---------------------------
    for heading, bullets in summary_dict.items():

        # combine bullet list to one block
        text_block = "\n".join(bullets)

        chunks = split_text_smartly(text_block, SLIDE_CHAR_LIMIT)

        # create 1+ slides depending on text size
        for i, chunk in enumerate(chunks):

            slide = copy_slide(prs, summary_template)
            clear_all_text(slide)
            add_watermark(slide)

            title_shape = get_title_shape(slide)
            body_shape = get_body_shape(slide)

            # Set heading
            title_shape.text = heading if i == 0 else f"{heading} (continued)"
            title_par = title_shape.text_frame.paragraphs[0]
            title_par.font.size = Pt(28)
            title_par.font.bold = True
            title_par.font.color.rgb = RGBColor(0, 0, 125)

            # Fill body
            
            tf = body_shape.text_frame
            tf.clear()

            for line in chunk.split("\n"):
                bullet = line.strip().lstrip("- ").strip()
                if not bullet:
                    continue

                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0, 0, 0) 
                p.level = 0
                p.space_after = Pt(5)

    # ---------------------------
    # FINAL SLIDE (Slide 3)
    # ---------------------------
    final_template = prs.slides[2]
    final_slide = copy_slide(prs, final_template)
    clear_all_text(final_slide)
    add_watermark(final_slide)
    add_thank_you_slide(prs, "Thank You", layout_index=3)


    body_shape = get_body_shape(final_slide)
    tf = body_shape.text_frame
    tf.word_wrap = True
    tf.clear()

    wrapped = wrap_url(item["URL"])

    p = tf.add_paragraph()
    p.text = wrapped
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 128)

    p.font.underline = True
    
    

    # ---------------------------
    # SAVE
    # ---------------------------
    output_path = f"{OUTPUT_FOLDER}/{clean_filename(item['Company'])}.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")

